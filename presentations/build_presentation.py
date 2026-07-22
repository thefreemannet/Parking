#!/usr/bin/env python3
"""Build the CAI2840C research presentation (English, Joshua Price 15-min structure)."""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "presentations" / "Parking_Occupancy_Research_Presentation.pptx"
FIG = ROOT / "outputs" / "figures"
DOWNLOADS = Path.home() / "Downloads" / "Parking_Occupancy_Research_Presentation.pptx"

BG = RGBColor(0xF7, 0xF8, 0xFA)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1F, 0x2E)
MUTED = RGBColor(0x5B, 0x64, 0x75)
ACCENT = RGBColor(0x0B, 0x6E, 0x7A)
ACCENT_SOFT = RGBColor(0xD6, 0xEE, 0xF1)
LINE = RGBColor(0xD8, 0xDE, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOOD = RGBColor(0x1B, 0x7A, 0x4E)

TOTAL = 11


def set_run(run, size=20, bold=False, color=INK, font="Calibri"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bg(slide, color=BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp_tree = slide.shapes._spTree
    sp = shape._element
    sp_tree.remove(sp)
    sp_tree.insert(2, sp)


def accent_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def footer(slide, page):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.1), Inches(10), Inches(0.3))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = f"CAI2840C  ·  Parking Occupancy Detection  ·  {page}/{TOTAL}"
    set_run(r, 11, False, MUTED)
    box2 = slide.shapes.add_textbox(Inches(11.5), Inches(7.1), Inches(1.4), Inches(0.3))
    p2 = box2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = "Gimeno · Jauregui · Rudes"
    set_run(r2, 11, False, MUTED)


def title_block(slide, title, subtitle=None, y=0.35):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(12.2), Inches(0.7))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = title
    set_run(r, 32, True, INK)
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.6), Inches(y + 0.65), Inches(12.2), Inches(0.45))
        r2 = box2.text_frame.paragraphs[0].add_run()
        r2.text = subtitle
        set_run(r2, 16, False, MUTED)


def card(slide, left, top, width, height):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = CARD
    sh.line.color.rgb = LINE
    sh.line.width = Pt(1)
    try:
        sh.adjustments[0] = 0.08
    except Exception:
        pass
    return sh


def bullets(slide, left, top, width, height, items, size=18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = "•  " + item
        set_run(r, size, False, INK)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def blank():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    accent_bar(slide)
    return slide


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
OUT.parent.mkdir(parents=True, exist_ok=True)

# --- 1 Title ---
s = blank()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
band.fill.solid()
band.fill.fore_color.rgb = ACCENT
band.line.fill.background()

eyebrow = s.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.5), Inches(0.4))
r = eyebrow.text_frame.paragraphs[0].add_run()
r.text = "CAI2840C  ·  Introduction to Computer Vision  ·  Research Presentation"
set_run(r, 14, True, ACCENT)

t = s.shapes.add_textbox(Inches(0.9), Inches(2.15), Inches(11.5), Inches(1.6))
tf = t.text_frame
tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "Parking Space Detection Using\nComputer Vision and Deep Learning"
set_run(r, 36, True, INK)

st = s.shapes.add_textbox(Inches(0.9), Inches(4.0), Inches(11.5), Inches(0.7))
tf = st.text_frame
tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = (
    "How reliably can deep learning detect parking-space occupancy\n"
    "across diverse real-world environments?"
)
set_run(r, 18, False, MUTED)

auth = s.shapes.add_textbox(Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.9))
tf = auth.text_frame
r = tf.paragraphs[0].add_run()
r.text = "Joaquin Gimeno  ·  Fernando Jauregui  ·  Eliot Rudes"
set_run(r, 16, True, INK)
p2 = tf.add_paragraph()
r2 = p2.add_run()
r2.text = (
    "Datasets: PKLot & CNRPark-EXT  ·  Models: Baseline CNN, MobileNetV3, VGG16, ResNet50"
)
set_run(r2, 14, False, MUTED)
notes(
    s,
    "Open with the research question. Parking waste matters for traffic, energy, and city "
    "space. One camera can watch many spots—but only if models stay accurate under rain, "
    "glare, and shadows. Preview: we compare three transfer-learning CNNs plus a baseline "
    "on two public datasets.",
)

# --- 2 Research question ---
s = blank()
footer(s, 2)
title_block(s, "Research question", "Why accuracy alone is not enough for smart parking")
card(s, Inches(0.6), Inches(1.6), Inches(6.0), Inches(4.8))
bullets(
    s,
    Inches(0.85),
    Inches(1.85),
    Inches(5.5),
    Inches(4.3),
    [
        "Cameras can monitor many spaces at once",
        "CNNs classify cropped spots; detectors locate vehicles in full scenes",
        "Benchmark scores often ignore rain, glare, shadows, occlusion",
        "Smart-city routing needs accuracy and speed together",
    ],
    17,
)
card(s, Inches(6.9), Inches(1.6), Inches(5.8), Inches(4.8))
box = s.shapes.add_textbox(Inches(7.15), Inches(1.85), Inches(5.3), Inches(0.4))
r = box.text_frame.paragraphs[0].add_run()
r.text = "Central question"
set_run(r, 14, True, ACCENT)
box = s.shapes.add_textbox(Inches(7.15), Inches(2.3), Inches(5.3), Inches(2.2))
tf = box.text_frame
tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = (
    "How reliably can deep learning–based computer vision detect parking-space occupancy "
    "across diverse real-world environments—and what do failures imply for future system design?"
)
set_run(r, 16, False, INK)
box = s.shapes.add_textbox(Inches(7.15), Inches(4.7), Inches(5.3), Inches(1.3))
tf = box.text_frame
tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "Sub-questions"
set_run(r, 13, True, ACCENT)
p = tf.add_paragraph()
r = p.add_run()
r.text = (
    "• Condition robustness  ·  Cross-dataset generalization\n"
    "• Failure patterns → sensor fusion  ·  Accuracy–speed trade-off"
)
set_run(r, 14, False, MUTED)
notes(
    s,
    "Spend about two minutes here. Emphasize that high lab accuracy is not the same as "
    "reliability in rain, shadows, and odd camera angles. Tie to research questions RQ1–RQ4.",
)

# --- 3 Contribution ---
s = blank()
footer(s, 3)
title_block(s, "Contribution to the literature", "Three gaps we address with one comparative experiment")
gaps = [
    (
        "1  Environmental reliability",
        "Few studies systematically stress-test models under rain, glare, shadows, and low light.",
    ),
    (
        "2  Cross-dataset generalization",
        "Models trained on one parking layout often fail on another (PKLot ↔ CNRPark-EXT).",
    ),
    (
        "3  Limits of vision-only systems",
        "Failure analysis reveals when cameras need radar / ultrasonic / IoT fusion.",
    ),
]
for i, (h, body) in enumerate(gaps):
    top = Inches(1.55 + i * 1.6)
    card(s, Inches(0.6), top, Inches(12.1), Inches(1.4))
    box = s.shapes.add_textbox(Inches(0.95), top + Inches(0.25), Inches(11.4), Inches(0.35))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = h
    set_run(r, 18, True, ACCENT)
    box = s.shapes.add_textbox(Inches(0.95), top + Inches(0.65), Inches(11.4), Inches(0.55))
    tf = box.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = body
    set_run(r, 16, False, INK)
notes(
    s,
    "Keep literature brief. Our contribution is a reproducible comparison under shared "
    "splits, metrics, and condition breakdowns—not a new architecture.",
)

# --- 4 Data ---
s = blank()
footer(s, 4)
title_block(s, "Data", "Two public datasets chosen for weather and scene diversity")
card(s, Inches(0.6), Inches(1.55), Inches(5.9), Inches(4.7))
box = s.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(5.3), Inches(0.4))
r = box.text_frame.paragraphs[0].add_run()
r.text = "PKLot"
set_run(r, 22, True, ACCENT)
bullets(
    s,
    Inches(0.9),
    Inches(2.35),
    Inches(5.3),
    Inches(3.5),
    [
        "Segmented parking-space patches",
        "Multiple lots: PUC, UFPR04, UFPR05",
        "Labeled weather: Sunny, Cloudy, Rainy",
        "Supports condition-level evaluation",
    ],
    16,
)
card(s, Inches(6.8), Inches(1.55), Inches(5.9), Inches(4.7))
box = s.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.3), Inches(0.4))
r = box.text_frame.paragraphs[0].add_run()
r.text = "CNRPark-EXT"
set_run(r, 22, True, ACCENT)
bullets(
    s,
    Inches(7.1),
    Inches(2.35),
    Inches(5.3),
    Inches(3.5),
    [
        "Real camera scenes → free / busy patches",
        "Different viewpoint and lot design",
        "Tests cross-dataset generalization (H2)",
        "Inventory records scene, camera, weather",
    ],
    16,
)
notes(
    s,
    "PKLot gives weather labels; CNRPark-EXT gives a second parking environment. "
    "Splits are scene/camera-aware when possible (seed 42) to reduce near-duplicate leakage.",
)

# --- 5 Methods ---
s = blank()
footer(s, 5)
title_block(s, "Methods", "Same split, same metrics — fair comparison")
models = [
    ("Baseline CNN", "Simple 3-layer CNN\nfrom scratch"),
    ("MobileNetV3", "Transfer learning\nImageNet backbone"),
    ("VGG16", "Transfer learning\nImageNet backbone"),
    ("ResNet50", "Transfer learning\nImageNet backbone"),
]
for i, (name, desc) in enumerate(models):
    left = Inches(0.6 + i * 3.15)
    card(s, left, Inches(1.55), Inches(3.0), Inches(2.3))
    box = s.shapes.add_textbox(left + Inches(0.2), Inches(1.8), Inches(2.6), Inches(0.45))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = name
    set_run(r, 16, True, ACCENT)
    box = s.shapes.add_textbox(left + Inches(0.2), Inches(2.35), Inches(2.6), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = desc
    set_run(r, 14, False, INK)

card(s, Inches(0.6), Inches(4.1), Inches(12.1), Inches(2.3))
bullets(
    s,
    Inches(0.9),
    Inches(4.3),
    Inches(11.5),
    Inches(2.0),
    [
        "Train-only augmentation: brightness, small rotation, shifts, horizontal flip",
        "Resize 224×224 · Adam · early stopping · best validation checkpoint",
        "Metrics: accuracy, precision, recall, F1, macro-F1, confusion matrix, model size, ms/image & FPS",
        "Optional track: YOLOv8s / YOLO11n for full-scene detection (off until labels exist)",
    ],
    15,
)
notes(
    s,
    "Transfer learning freezes ImageNet features and trains a small head. Identical protocol "
    "so differences reflect architecture. YOLO is an optional secondary task.",
)

# --- 6 Pipeline ---
s = blank()
footer(s, 6)
title_block(s, "Experimental pipeline", "Reproducible end-to-end workflow")
steps = [
    ("01", "Inventory", "Scan PKLot &\nCNRPark-EXT labels"),
    ("02", "Split", "Scene-aware\n70 / 15 / 15"),
    ("03", "Train", "Baseline +\n3 TL models"),
    ("04", "Evaluate", "Test metrics &\ncondition tables"),
    ("05", "Review", "Failure cases\nfor H3"),
]
for i, (num, title, body) in enumerate(steps):
    left = Inches(0.55 + i * 2.5)
    card(s, left, Inches(1.7), Inches(2.3), Inches(3.5))
    box = s.shapes.add_textbox(left + Inches(0.15), Inches(1.95), Inches(2.0), Inches(0.35))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = num
    set_run(r, 14, True, ACCENT)
    box = s.shapes.add_textbox(left + Inches(0.15), Inches(2.4), Inches(2.0), Inches(0.4))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = title
    set_run(r, 18, True, INK)
    box = s.shapes.add_textbox(left + Inches(0.15), Inches(3.0), Inches(2.0), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = body
    set_run(r, 14, False, MUTED)
    if i < 4:
        arrow = s.shapes.add_textbox(left + Inches(2.15), Inches(3.1), Inches(0.4), Inches(0.4))
        r = arrow.text_frame.paragraphs[0].add_run()
        r.text = "→"
        set_run(r, 22, True, ACCENT)

box = s.shapes.add_textbox(Inches(0.6), Inches(5.6), Inches(12.1), Inches(0.9))
tf = box.text_frame
tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = (
    "Code & config archived: configs/default.yaml · "
    "notebooks/01_parking_occupancy_experiment.ipynb · outputs/metrics/"
)
set_run(r, 14, False, MUTED)
notes(
    s,
    "Walk the pipeline quickly. Stress reproducibility: fixed seed, shared config, "
    "archived split manifest and metrics JSON.",
)

# --- 7 Results table ---
s = blank()
footer(s, 7)
title_block(
    s,
    "Results — classification comparison",
    "Preliminary demo run on the reproducible pipeline (synthetic patches for pipeline validation)",
)
headers = ["Model", "Accuracy", "F1", "Size (MB)", "ms / image", "FPS"]
rows = [
    ["Baseline CNN", "0.50", "0.00", "1.1", "1.7", "583"],
    ["MobileNetV3", "0.83", "0.85", "4.1", "3.7", "268"],
    ["VGG16*", "0.00", "0.00", "56.2", "38.7", "26"],
    ["ResNet50", "0.81", "0.76", "90.6", "22.4", "45"],
]
col_w = [2.4, 1.6, 1.4, 1.8, 2.0, 1.4]
left0 = Inches(1.2)
top0 = Inches(1.7)
x = left0
for j, h in enumerate(headers):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top0, Inches(col_w[j]), Inches(0.55))
    sh.fill.solid()
    sh.fill.fore_color.rgb = ACCENT
    sh.line.fill.background()
    box = s.shapes.add_textbox(x, top0 + Inches(0.1), Inches(col_w[j]), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = h
    set_run(r, 13, True, WHITE)
    x += Inches(col_w[j])

for i, row in enumerate(rows):
    y = top0 + Inches(0.55 + i * 0.6)
    x = left0
    bgc = CARD if i % 2 == 0 else ACCENT_SOFT
    for j, val in enumerate(row):
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(col_w[j]), Inches(0.6))
        sh.fill.solid()
        sh.fill.fore_color.rgb = bgc
        sh.line.color.rgb = LINE
        sh.line.width = Pt(0.5)
        box = s.shapes.add_textbox(x, y + Inches(0.12), Inches(col_w[j]), Inches(0.4))
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = val
        set_run(r, 14, i == 1, GOOD if i == 1 else INK)
        x += Inches(col_w[j])

note = s.shapes.add_textbox(Inches(1.2), Inches(5.5), Inches(10.8), Inches(1.0))
tf = note.text_frame
tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = (
    "Takeaway: MobileNetV3 offers the best accuracy–speed–size balance in this run "
    "(supports H4). *VGG16 collapsed on synthetic demo data—full PKLot/CNRPark-EXT "
    "training is required for research claims."
)
set_run(r, 14, False, MUTED)
notes(
    s,
    "Do not just read the table. Interpret: MobileNetV3 wins the practical trade-off. "
    "ResNet50 is competitive but heavier. Flag that these numbers validate the pipeline; "
    "full-dataset results are the research claims.",
)

# --- 8 Confusion matrices ---
s = blank()
footer(s, 8)
title_block(
    s,
    "Results — confusion matrices",
    "Where each model is right and wrong on the held-out test set",
)
imgs = [
    ("MobileNetV3", FIG / "mobilenetv3_confusion_matrix.png"),
    ("ResNet50", FIG / "resnet50_confusion_matrix.png"),
    ("Baseline", FIG / "baseline_confusion_matrix.png"),
]
for i, (label, path) in enumerate(imgs):
    left = Inches(0.5 + i * 4.2)
    box = s.shapes.add_textbox(left, Inches(1.45), Inches(4.0), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    set_run(r, 14, True, ACCENT)
    if path.exists():
        s.shapes.add_picture(str(path), left + Inches(0.25), Inches(1.9), height=Inches(4.5))
notes(
    s,
    "Walk through one matrix: true empty vs occupied. MobileNetV3 recovers both classes; "
    "baseline collapses toward one class.",
)

# --- 9 Conditions ---
s = blank()
footer(s, 9)
title_block(s, "Conditions & failure patterns", "Connecting results to H1–H3")
card(s, Inches(0.6), Inches(1.55), Inches(6.0), Inches(4.8))
box = s.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(5.4), Inches(0.4))
r = box.text_frame.paragraphs[0].add_run()
r.text = "MobileNetV3 by weather (demo)"
set_run(r, 16, True, ACCENT)
conds = [("Rainy", "100%", 1.0), ("Cloudy", "89%", 0.89), ("Sunny", "75%", 0.75)]
for i, (name, pct, v) in enumerate(conds):
    y = Inches(2.45 + i * 1.0)
    box = s.shapes.add_textbox(Inches(0.9), y, Inches(1.6), Inches(0.35))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = name
    set_run(r, 14, False, INK)
    track = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.6), y + Inches(0.05), Inches(3.4), Inches(0.32)
    )
    track.fill.solid()
    track.fill.fore_color.rgb = LINE
    track.line.fill.background()
    bar = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.6), y + Inches(0.05), Inches(3.4 * v), Inches(0.32)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    box = s.shapes.add_textbox(Inches(5.1), y, Inches(1.0), Inches(0.35))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = pct
    set_run(r, 14, True, INK)
box = s.shapes.add_textbox(Inches(0.9), Inches(5.5), Inches(5.4), Inches(0.6))
tf = box.text_frame
tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "Also tracked: dataset, camera, and scene accuracy when n ≥ 5."
set_run(r, 13, False, MUTED)

card(s, Inches(6.9), Inches(1.55), Inches(5.8), Inches(4.8))
box = s.shapes.add_textbox(Inches(7.2), Inches(1.8), Inches(5.3), Inches(0.4))
r = box.text_frame.paragraphs[0].add_run()
r.text = "What failures suggest"
set_run(r, 16, True, ACCENT)
bullets(
    s,
    Inches(7.2),
    Inches(2.35),
    Inches(5.2),
    Inches(3.6),
    [
        "Vision-only models still struggle at the edges of lighting and layout",
        "Cross-dataset gaps support the need for diverse training scenes",
        "Exported failure CSVs support qualitative review (shadows, glare, occlusion)",
        "Points toward future sensor-fusion—not replacing cameras, complementing them",
    ],
    15,
)
notes(
    s,
    "Link to hypotheses: H1 condition drops, H2 dataset differences, H3 failure cases "
    "justify fusion research. Be honest that demo weather labels are synthetic.",
)

# --- 10 Conclusion ---
s = blank()
footer(s, 10)
title_block(s, "What the answer means", "Implications for real-time smart parking")
points = [
    (
        "Best practical balance",
        "MobileNetV3 combines strong F1 with small size and high FPS—best fit for "
        "edge/real-time deployment among the four.",
    ),
    (
        "Reliability ≠ one number",
        "Condition and dataset breakdowns matter as much as overall accuracy for city systems.",
    ),
    (
        "Vision has limits",
        "Failure patterns under hard conditions motivate camera + ultrasonic / radar / IoT "
        "fusion research.",
    ),
    (
        "Reproducible framework",
        "Shared splits, configs, and metrics let others extend the study to full "
        "PKLot/CNRPark-EXT and optional YOLO.",
    ),
]
for i, (h, body) in enumerate(points):
    top = Inches(1.5 + i * 1.2)
    card(s, Inches(0.6), top, Inches(12.1), Inches(1.05))
    box = s.shapes.add_textbox(Inches(0.95), top + Inches(0.15), Inches(11.4), Inches(0.3))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = h
    set_run(r, 16, True, ACCENT)
    box = s.shapes.add_textbox(Inches(0.95), top + Inches(0.5), Inches(11.4), Inches(0.4))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = body
    set_run(r, 14, False, INK)
notes(
    s,
    "Conclude—do not just summarize. Prefer lightweight models for cameras; report "
    "condition metrics; plan fusion where vision fails.",
)

# --- 11 Thank you ---
s = blank()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
band.fill.solid()
band.fill.fore_color.rgb = ACCENT
band.line.fill.background()

box = s.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.6))
r = box.text_frame.paragraphs[0].add_run()
r.text = "Thank you"
set_run(r, 40, True, INK)

box = s.shapes.add_textbox(Inches(0.9), Inches(2.9), Inches(11.5), Inches(0.8))
tf = box.text_frame
tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "Questions?\nRepo: github.com/thefreemannet/Parking"
set_run(r, 20, False, MUTED)

box = s.shapes.add_textbox(Inches(0.9), Inches(4.8), Inches(11.5), Inches(1.2))
tf = box.text_frame
tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = (
    "Next steps: full PKLot + CNRPark-EXT training · deeper fine-tuning · "
    "optional YOLOv8s/YOLO11n full-scene track"
)
set_run(r, 15, False, ACCENT)

box = s.shapes.add_textbox(Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.4))
r = box.text_frame.paragraphs[0].add_run()
r.text = "Joaquin Gimeno  ·  Fernando Jauregui  ·  Eliot Rudes  ·  CAI2840C"
set_run(r, 14, False, MUTED)
notes(
    s,
    "Invite questions. Have failure CSV and confusion matrices ready. Mention GitHub for "
    "reproducibility.",
)

prs.save(str(OUT))
shutil.copy2(OUT, DOWNLOADS)
print(f"Saved: {OUT}")
print(f"Copy:  {DOWNLOADS}")
print(f"Slides: {len(prs.slides)}")
